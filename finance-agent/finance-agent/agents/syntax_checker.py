"""Syntax Checker Agent — validates generated code before execution.

Two-phase checking:
1. Python AST parsing for basic syntax errors (fast, local)
2. LLM-based deep API review for library-specific mistakes (thorough)
"""

import ast
import json
import re
import logging
import time
import httpx
from core.config import config

log = logging.getLogger("syntax_checker")

MAX_RETRIES = 2  # Max retry attempts for LLM calls (total attempts = MAX_RETRIES + 1)
MAX_CODE_LENGTH = 15000

SYNTAX_CHECKER_SYSTEM_PROMPT = """You are an expert Python syntax and API reviewer specializing in document-generation and data-visualization libraries. Your ONLY job is to find syntax errors, API misuse, and incorrect method calls in Python code.

You have DEEP knowledge of these libraries and their exact APIs:

═══════════════════════════════════════════════
PYTHON-PPTX (pptx) - Presentation Generation
═══════════════════════════════════════════════
CORRECT IMPORTS:
  from pptx import Presentation
  from pptx.util import Inches, Pt, Emu, Cm
  from pptx.dml.color import RGBColor
  from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
  from pptx.enum.chart import XL_CHART_TYPE

COMMON MISTAKES TO CATCH:
  ✗ RGBColor('#1a73e8')          → ✓ RGBColor(0x1a, 0x73, 0xe8) [takes 3 ints, NOT hex string]
  ✗ RGBColor('1a73e8')           → ✓ RGBColor(0x1a, 0x73, 0xe8)
  ✗ paragraph.font.size = 14     → ✓ paragraph.font.size = Pt(14)
  ✗ paragraph.font.color = ...   → ✓ paragraph.font.color.rgb = RGBColor(...)
  ✗ shape.left = 1               → ✓ shape.left = Inches(1)
  ✗ slide_layout = prs.slide_layouts['Title'] → ✓ prs.slide_layouts[0] [index, not string]
  ✗ prs.add_slide(layout)        → ✓ prs.slides.add_slide(layout)
  ✗ from pptx.util import Points → ✓ from pptx.util import Pt
  ✗ paragraph.alignment = 'center' → ✓ paragraph.alignment = PP_ALIGN.CENTER
  ✗ table.cell(row, col).text = number → must convert to str first
  ✗ slide.shapes.add_table(rows, cols) → missing position/size args

═══════════════════════════════════════════════
REPORTLAB - PDF Generation
═══════════════════════════════════════════════
CORRECT IMPORTS:
  from reportlab.lib.pagesizes import letter, A4
  from reportlab.lib.units import inch, cm, mm
  from reportlab.lib.colors import HexColor, Color, black, white, blue
  from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
  from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
  from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak

COMMON MISTAKES TO CATCH:
  ✗ HexColor('1a73e8')           → ✓ HexColor('#1a73e8') [NEEDS # prefix]
  ✗ from reportlab.lib.colors import RGBColor → ✓ HexColor or Color(r,g,b)
  ✗ doc = SimpleDocTemplate('file.pdf') → needs pagesize: SimpleDocTemplate('file.pdf', pagesize=letter)
  ✗ Paragraph(text)              → ✓ Paragraph(text, style) [style is required]
  ✗ Spacer(12)                   → ✓ Spacer(1, 12) [needs width and height]
  ✗ Table(data, colWidths=100)   → ✓ Table(data, colWidths=[100]*n) [must be list]
  ✗ Image('path', 400, 300)      → ✓ Image('path', width=400, height=300)
  ✗ elements.append(text_string) → ✓ elements.append(Paragraph(text_string, style))

═══════════════════════════════════════════════
PYTHON-DOCX (docx) - Word Document Generation
═══════════════════════════════════════════════
COMMON MISTAKES TO CATCH:
  ✗ doc = Document('new.docx')   → ✓ doc = Document() [no arg for new doc]
  ✗ run.font.color = RGBColor(...)  → ✓ run.font.color.rgb = RGBColor(...)
  ✗ RGBColor('#1a73e8')          → ✓ RGBColor(0x1a, 0x73, 0xe8) [3 ints]
  ✗ from docx.shared import Points → ✓ from docx.shared import Pt
  ✗ doc.add_picture(path, Inches(4)) → ✓ doc.add_picture(path, width=Inches(4))
  ✗ paragraph.alignment = 'center' → ✓ paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
  ✗ doc.save() → ✓ doc.save('filename.docx') [path required]

═══════════════════════════════════════════════
MATPLOTLIB - Chart Generation
═══════════════════════════════════════════════
COMMON MISTAKES TO CATCH:
  ✗ plt.show() → should use plt.savefig() then plt.close() in headless mode
  ✗ Missing matplotlib.use('Agg') before creating figures
  ✗ plt.savefig('chart.png') without using output_dir
  ✗ fig.savefig() without plt.close() → memory leak
  ✗ Defining a variable named `colors` — shadows reportlab import

═══════════════════════════════════════════════
GENERAL PYTHON MISTAKES TO CATCH:
═══════════════════════════════════════════════
  ✗ f-strings with unescaped braces
  ✗ Missing imports
  ✗ String/number type mismatches
  ✗ Using variables before assignment
  ✗ Unclosed parentheses/brackets/strings
  ✗ Division by zero possibilities

YOUR RESPONSE FORMAT:
{
  "has_errors": true/false,
  "errors": [
    {
      "line": "approximate line number or code snippet",
      "error_type": "syntax|api_misuse|import|type_mismatch|missing_import|logic",
      "description": "What is wrong",
      "fix": "The corrected code for that section",
      "severity": "critical|warning"
    }
  ],
  "fixed_code": "the complete corrected Python code if errors were found, or empty string if no errors",
  "summary": "Brief summary of findings"
}

IMPORTANT RULES:
1. ONLY report REAL errors. Do not flag correct code as wrong.
2. If the code is correct, return has_errors: false with empty errors array.
3. For critical errors, provide the fixed_code with ALL corrections applied.
4. Preserve the original code's logic and intent — only fix errors, don't refactor.
5. Pay special attention to RGBColor usage — this is the #1 most common mistake.
6. Check every import statement against the correct import paths listed above."""


async def _call_llm(system_prompt: str, user_message: str, label: str = "SYNTAX") -> str:
    """Call OpenRouter for syntax checking."""
    last_error = None
    for attempt in range(MAX_RETRIES + 1):
        try:
            log.info("[%s] Calling OpenRouter (attempt=%d/%d)...", label, attempt + 1, MAX_RETRIES + 1)
            t0 = time.time()
            async with httpx.AsyncClient(timeout=120) as client:
                response = await client.post(
                    "https://openrouter.ai/api/v1/chat/completions",
                    headers={
                        "Authorization": f"Bearer {config.openrouter_api_key}",
                        "Content-Type": "application/json",
                    },
                    json={
                        "model": config.openrouter_model,
                        "messages": [
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": user_message},
                        ],
                        "max_tokens": 8192,
                        "temperature": 0.1,
                        "top_p": 0.95,
                        "reasoning": {"effort": "high"},
                    },
                )
                elapsed = time.time() - t0
                response.raise_for_status()
                data = response.json()
                content = data["choices"][0]["message"]["content"]

                # Strip <think> blocks
                content = re.sub(r"<think>.*?</think>", "", content, flags=re.DOTALL).strip()
                # Strip markdown fences
                if content.startswith("```"):
                    content = content.split("\n", 1)[1] if "\n" in content else content[3:]
                if content.endswith("```"):
                    content = content[:-3]
                # Strip pre-JSON text
                content = content.strip()
                if content and not content.startswith("{"):
                    idx = content.find("{")
                    if idx > 0:
                        content = content[idx:]

                log.info("[%s] Response in %.1fs", label, elapsed)
                return content.strip()
        except (httpx.HTTPStatusError, httpx.ConnectError, httpx.ReadTimeout) as e:
            last_error = e
            log.warning("[%s] Attempt %d failed: %s", label, attempt + 1, e)
            if attempt < MAX_RETRIES:
                import asyncio
                await asyncio.sleep(2)
    raise last_error


def _parse_json(raw: str) -> dict:
    """Parse JSON from LLM response."""
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        pass
    match = re.search(r"\{.*\}", raw, re.DOTALL)
    if match:
        try:
            return json.loads(match.group())
        except json.JSONDecodeError:
            pass
    return {"has_errors": False, "errors": [], "fixed_code": "", "summary": "LLM response unparseable — proceeding with original code"}


def check_ast(code: str) -> dict:
    """Phase 1: Quick Python AST parse to catch basic syntax errors."""
    try:
        ast.parse(code)
        return {"has_errors": False, "errors": [], "fixed_code": "", "summary": "AST parse OK"}
    except SyntaxError as e:
        return {
            "has_errors": True,
            "errors": [{
                "line": str(e.lineno or "unknown"),
                "error_type": "syntax",
                "description": f"SyntaxError: {e.msg} (line {e.lineno}, col {e.offset})",
                "fix": "",
                "severity": "critical",
            }],
            "fixed_code": "",
            "summary": f"Python syntax error at line {e.lineno}: {e.msg}",
        }


async def check(code: str) -> dict:
    """Check code for syntax errors and API misuse.

    Phase 1: AST parse (fast, local).
    Phase 2: LLM deep API review (catches library-specific mistakes).

    Returns dict with: has_errors, errors, fixed_code, summary
    """
    if not code or not code.strip():
        return {
            "has_errors": True,
            "errors": [{"line": "0", "error_type": "syntax", "description": "Empty code", "fix": "", "severity": "critical"}],
            "fixed_code": "",
            "summary": "Code is empty",
        }

    # Phase 1: Quick AST syntax check
    ast_result = check_ast(code)
    if ast_result["has_errors"]:
        log.info("[SYNTAX] AST error found: %s", ast_result["errors"][0]["description"])
        # Ask LLM to fix syntax errors
        error_desc = "\n".join(f"- {e['description']}" for e in ast_result["errors"])
        truncated = code[:MAX_CODE_LENGTH]
        raw = await _call_llm(
            SYNTAX_CHECKER_SYSTEM_PROMPT,
            f"This Python code has syntax errors detected by the AST parser:\n\nERRORS:\n{error_desc}\n\nCODE:\n```python\n{truncated}\n```\n\nFix ALL syntax errors and also check for any API misuse in document-generation libraries.",
            label="SYNTAX-FIX",
        )
        result = _parse_json(raw)
        result.setdefault("has_errors", True)
        result.setdefault("errors", ast_result["errors"])
        result.setdefault("fixed_code", "")
        result.setdefault("summary", "Syntax errors found")
        return result

    # Phase 2: Deep LLM-based API review
    truncated = code[:MAX_CODE_LENGTH] if len(code) > MAX_CODE_LENGTH else code
    log.info("[SYNTAX] AST OK — running deep LLM API review (%d chars)", len(truncated))
    raw = await _call_llm(
        SYNTAX_CHECKER_SYSTEM_PROMPT,
        f"Review this Python code for syntax errors and API misuse. Focus especially on python-pptx, reportlab, python-docx, and matplotlib APIs.\n\nCODE:\n```python\n{truncated}\n```",
        label="SYNTAX-REVIEW",
    )
    result = _parse_json(raw)
    result.setdefault("has_errors", False)
    result.setdefault("errors", [])
    result.setdefault("fixed_code", "")
    result.setdefault("summary", "No issues found")
    return result
